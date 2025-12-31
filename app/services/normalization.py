from __future__ import annotations
import os, tempfile, time, html, io
from typing import Any, Dict, List, Optional, Tuple
from bs4 import BeautifulSoup, Tag

# Optional deps; guard imports to make normalization robust across environments
try:
    from paddleocr import PaddleOCR  # type: ignore
    _HAS_PADDLE = True
except Exception:
    PaddleOCR = None  # type: ignore
    _HAS_PADDLE = False
try:
    import pytesseract  # type: ignore
    _HAS_TESS = True
except Exception:
    pytesseract = None  # type: ignore
    _HAS_TESS = False
try:
    from PIL import Image  # type: ignore
    _HAS_PIL = True
except Exception:
    Image = None  # type: ignore
    _HAS_PIL = False
try:
    import fitz  # type: ignore  # PyMuPDF
    _HAS_FITZ = True
except Exception:
    fitz = None  # type: ignore
    _HAS_FITZ = False
try:
    import mammoth  # type: ignore
    _HAS_MAMMOTH = True
except Exception:
    mammoth = None  # type: ignore
    _HAS_MAMMOTH = False
try:
    # Optional: unstructured for native PDF text extraction
    from unstructured.partition.pdf import partition_pdf  # type: ignore
    _HAS_UNSTRUCTURED = True
except Exception:
    partition_pdf = None  # type: ignore
    _HAS_UNSTRUCTURED = False
try:
    import pdfplumber  # type: ignore
    _HAS_PDFPLUMBER = True
except Exception:
    pdfplumber = None  # type: ignore
    _HAS_PDFPLUMBER = False
try:
    from pdfminer.high_level import extract_text as pdfminer_extract_text  # type: ignore
    _HAS_PDFMINER = True
except Exception:
    pdfminer_extract_text = None  # type: ignore
    _HAS_PDFMINER = False
try:
    from pptx import Presentation  # type: ignore
    _HAS_PPTX = True
except Exception:
    Presentation = None  # type: ignore
    _HAS_PPTX = False
try:
    import openpyxl  # type: ignore
    _HAS_OPENPYXL = True
except Exception:
    openpyxl = None  # type: ignore
    _HAS_OPENPYXL = False
try:
    from langdetect import detect as lang_detect  # type: ignore
    _HAS_LANG = True
except Exception:
    lang_detect = None  # type: ignore
    _HAS_LANG = False

from bs4 import BeautifulSoup  # already imported
from infra.db import DBClient
from infra.minio_store import MinioStore
from core.config import settings
from .manifests import CanonicalArtifact, CanonicalManifest, new_artifact_id
from .parsing import DocumentParserManager, AdvancedParserAdapter, SimpleFallbackAdapter, ParserOptions


def sanitize_html(raw_html: str) -> str:
    soup = BeautifulSoup(raw_html, "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    # Ensure one <body>
    if not soup.body:
        body = soup.new_tag("body")
        body.append(BeautifulSoup(raw_html, "lxml"))
        soup.append(body)
    return str(soup)


def wrap_txt_to_html(text: str) -> str:
    body = f"<pre>{html.escape(text)}</pre>"
    return f"<!doctype html><html><head><meta charset='utf-8'></head><body>{body}</body></html>"


_POCR = None  # lazy init PaddleOCR
def _ocr_image_to_text(img) -> tuple[str, list[str]]:
    """OCR wrapper using PaddleOCR only. Returns (text, warnings)."""
    warnings: list[str] = []
    if img is None or not _HAS_PIL:
        missing = []
        if not _HAS_PIL: missing.append("PIL")
        warnings.append(f"ocr_tool_missing:{'+'.join(missing) if missing else 'unknown'}")
        return "", warnings

    if _HAS_PADDLE:
        try:
            global _POCR
            if _POCR is None:
                # Map common codes to Paddle lang
                langs = settings.ocr_langs.lower()
                lang = "en" if "en" in langs or "eng" in langs else "en"
                _POCR = PaddleOCR(lang=lang, use_angle_cls=True, use_gpu=False, show_log=False)  # type: ignore
            # Write image to a temp PNG to avoid numpy dependency
            import tempfile
            tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            tmp_path = tmp.name
            tmp.close()
            try:
                img.save(tmp_path, format="PNG")
                result = _POCR.ocr(tmp_path, cls=True)  # type: ignore
            finally:
                try:
                    os.remove(tmp_path)
                except Exception:
                    pass
            lines: list[str] = []
            for det in (result or []):
                if not isinstance(det, list):
                    continue
                for entry in det:
                    try:
                        # entry: [box, (text, conf)]
                        text = (entry[1][0] if isinstance(entry[1], (list, tuple)) else None) or ""
                        if text:
                            lines.append(str(text))
                    except Exception:
                        continue
            text_out = "\n".join(lines).strip()
            return text_out, warnings
        except Exception as e:
            warnings.append(f"paddle_ocr_failed:{e}")
            return "", warnings

    # Fallback: Tesseract OCR if available
    if _HAS_TESS and _HAS_PIL and img is not None:
        try:
            txt = pytesseract.image_to_string(img)  # type: ignore
            return (txt or "").strip(), warnings
        except Exception as e:
            warnings.append(f"tesseract_failed:{e}")

    warnings.append("ocr_tool_missing:paddleocr")
    return "", warnings


def pdf_to_html(tmp_path: str) -> Tuple[str, int, int, list[str]]:
    """PDF pipeline with quick classification per page.
    - If a page has a usable text layer, extract with MuPDF (or unstructured when available for entire-doc fast parse).
    - If a page lacks text, render to image and OCR (PaddleOCR).
    Returns: (html, page_count, ocr_pages, warnings)
    """
    warnings: list[str] = []
    if not _HAS_FITZ:
        warnings.append("pdf_tool_missing:pymupdf")
        # Try pdfminer as a plain-text fallback
        if _HAS_PDFMINER:
            try:
                text = pdfminer_extract_text(tmp_path) or ""
                html_doc = wrap_txt_to_html(text)
                return html_doc, 1, 0, warnings
            except Exception as e:
                warnings.append(f"pdfminer_failed:{e}")
        return wrap_txt_to_html("PDF content (no extractor available)"), 1, 0, warnings

    doc = fitz.open(tmp_path)  # type: ignore
    # caps - Quality Over Speed: Effectively disable limits
    OCR_MAX_PAGES = 10000 
    PDF_NATIVE_ONLY_IF_PAGES_GT = 10000

    # Quick doc-level detection: do we have any text at all?
    total_text_chars = 0
    for page in doc:
        try:
            total_text_chars += len((page.get_text("text") or "").strip())
        except Exception:
            pass
    has_text_layer = total_text_chars > 100  # small threshold

    # If entirely native and unstructured is available, let it extract a plain text rendition once
    unstructured_text: Optional[str] = None
    if has_text_layer and _HAS_UNSTRUCTURED:
        try:
            elements = partition_pdf(filename=tmp_path, strategy="fast")  # type: ignore
            texts = [getattr(el, "text", "") for el in elements if getattr(el, "text", None)]
            unstructured_text = "\n\n".join(texts)
            warnings.append("native_pdf_unstructured")
        except Exception as e:
            warnings.append(f"unstructured_failed:{e}")
            unstructured_text = None

    ocr_pages = 0
    pages_html: list[str] = []
    plumber_tables_total = 0
    # If very large doc: do not attempt OCR; rely on text layers only
    large_native_only = doc.page_count > PDF_NATIVE_ONLY_IF_PAGES_GT
    if large_native_only:
        warnings.append("pdf_native_only_due_to_size")
    plumber_doc = None
    if _HAS_PDFPLUMBER:
        try:
            plumber_doc = pdfplumber.open(tmp_path)  # type: ignore
        except Exception as e:
            warnings.append(f"pdfplumber_failed:{e}")
            plumber_doc = None

    try:
        for i, page in enumerate(doc):
            page_num = i + 1
            # 1. Attempt Table Extraction & De-duplication Setup
            plumber_tables_found = []
            if plumber_doc:
                try:
                    if page_num - 1 < len(plumber_doc.pages):
                        p_plumber = plumber_doc.pages[page_num - 1]
                        # Strategy 1: Default
                        found_objs = p_plumber.find_tables() or []
                        # Strategy 2: Text alignment (if no tables found)
                        if not found_objs:
                            try:
                                found_objs = p_plumber.find_tables(table_settings={
                                    "vertical_strategy": "text", 
                                    "horizontal_strategy": "text"
                                }) or []
                            except Exception:
                                pass
                        
                        if found_objs:
                            plumber_tables_total += len(found_objs)
                            for tbl in found_objs:
                                # A. Redact from PyMuPDF Page (De-duplication)
                                if hasattr(tbl, 'bbox'):
                                    # bbox: (x0, top, x1, bottom)
                                    page.add_redact_annot(fitz.Rect(tbl.bbox))
                                
                                # B. Extract Data
                                data = tbl.extract()
                                if data:
                                    rows = [r for r in data if any((c or '').strip() for c in r)]
                                    if rows:
                                        tr_html = []
                                        for row in rows:
                                            cells = ''.join(f"<td>{html.escape(str(c or ''))}</td>" for c in row)
                                            tr_html.append(f"<tr>{cells}</tr>")
                                        plumber_tables_found.append(f"<div data-source='pdfplumber'><table>{''.join(tr_html)}</table></div>")
                except Exception as e:
                    warnings.append(f"pdfplumber_page_setup_failed:{page_num}:{e}")

            try:
                page_text = (page.get_text("text") or "").strip()
            except Exception:
                page_text = ""

            # Quality Check: Trust native if content > 50 chars OR if we successfully extracted tables
            # If tables found, we trust the PDF is native enough to use.
            if len(page_text) > 50 or len(plumber_tables_found) > 0:
                # Native text layer present
                
                # Apply Redactions to remove table text from the layout analysis
                if plumber_tables_found:
                    try:
                        # apply_redactions removes the content covered by the annotations
                        page.apply_redactions()
                    except Exception:
                        pass

                try:
                    frag = page.get_text("xhtml")  # type: ignore
                    # Strip the xml declaration and body tags
                    if "<body>" in frag:
                        frag = frag.split("<body>")[1].split("</body>")[0]
                except Exception:
                    # If xhtml fails (or after redaction something weird happens)
                    frag = html.escape((page.get_text("text") or "").strip())
                    frag = f"<pre>{frag}</pre>"

                # Append the clean structured tables
                if plumber_tables_found:
                    frag += "".join(plumber_tables_found)
                
                pages_html.append(f"<section data-page='{page_num}'>{frag}</section>")
                continue

            # Scanned page: rasterize and OCR (subject to caps)
            if large_native_only:
                # keep placeholder empty section to preserve page numbering
                pages_html.append(f"<section data-page='{page_num}'><pre></pre></section>")
                continue
            if ocr_pages >= OCR_MAX_PAGES:
                warnings.append("ocr_skipped_due_to_cap")
                pages_html.append(f"<section data-page='{page_num}'><pre></pre></section>")
                continue
            try:
                # 3x scale for Maximum OCR quality (approx 300 DPI)
                pix = page.get_pixmap(matrix=fitz.Matrix(3, 3))  # type: ignore
                if _HAS_PIL:
                    img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")  # type: ignore
                else:
                    img = None
            except Exception as e:
                warnings.append(f"render_failed:{e}")
                img = None
            text, w = _ocr_image_to_text(img)
            warnings.extend(w)
            if text:
                ocr_pages += 1
            pages_html.append(f"<section data-page='{page_num}'><pre>{html.escape(text)}</pre></section>")

    finally:
        if plumber_doc:
            try:
                plumber_doc.close()
            except Exception:
                pass

    html_doc = f"<!doctype html><html><head><meta charset='utf-8'></head><body>{''.join(pages_html)}</body></html>"

    # If unstructured text is available and the MuPDF text is sparse, include it as a hidden appendix for recall
    if unstructured_text and has_text_layer:
        appendix = f"<section data-page='0' hidden><pre>{html.escape(unstructured_text)}</pre></section>"
        html_doc = html_doc.replace("</body>", appendix + "</body>")

    text_content = BeautifulSoup(html_doc, 'lxml').get_text() or ""
    if len(text_content.strip()) == 0:
        warnings.append("canonical_empty")
    if ocr_pages > 0:
        warnings.append(f"ocr_pages:{ocr_pages}")
    if plumber_tables_total > 0:
        warnings.append(f"plumber_tables:{plumber_tables_total}")
    return html_doc, doc.page_count, ocr_pages, warnings


def docx_to_html(tmp_path: str) -> Tuple[str, int, int, list[str]]:
    if not _HAS_MAMMOTH:
        return wrap_txt_to_html("DOCX conversion not available"), 1, 0, ["docx_tool_missing:mammoth"]
    with open(tmp_path, "rb") as f:
        result = mammoth.convert_to_html(f)  # type: ignore
    html_doc = f"<!doctype html><html><head><meta charset='utf-8'></head><body>{result.value}</body></html>"
    return sanitize_html(html_doc), 1, 0, []


def pptx_to_html(tmp_path: str) -> Tuple[str, int, int, list[str]]:
    warnings: list[str] = []
    if not _HAS_PPTX:
        warnings.append("pptx_tool_missing:python-pptx")
        return wrap_txt_to_html("PPTX conversion not available"), 1, 0, warnings
    try:
        prs = Presentation(tmp_path)  # type: ignore
        slides = []
        for idx, s in enumerate(prs.slides, start=1):
            texts = []
            for shape in s.shapes:
                try:
                    if hasattr(shape, "text"):
                        t = (shape.text or "").strip()
                        if t:
                            texts.append(t)
                except Exception:
                    continue
            content = html.escape("\n\n".join(texts))
            slides.append(f"<section data-page='{idx}'><pre>{content}</pre></section>")
        html_doc = f"<!doctype html><html><head><meta charset='utf-8'></head><body>{''.join(slides)}</body></html>"
        return sanitize_html(html_doc), len(slides) or 1, 0, warnings
    except Exception as e:
        warnings.append(f"pptx_failed:{e}")
        return wrap_txt_to_html("PPTX parse failed"), 1, 0, warnings


def xlsx_to_html(tmp_path: str) -> Tuple[str, int, int, list[str]]:
    warnings: list[str] = []
    if not _HAS_OPENPYXL:
        warnings.append("xlsx_tool_missing:openpyxl")
        return wrap_txt_to_html("XLSX conversion not available"), 1, 0, warnings
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)  # type: ignore
        sheets_html = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [html.escape(str(c) if c is not None else "") for c in row]
                rows.append("<tr>" + "".join(f"<td>{c}</td>" for c in cells) + "</tr>")
            table = f"<table>{''.join(rows)}</table>"
            sheets_html.append(f"<section data-page='1'><div data-sheet='{html.escape(ws.title)}'>{table}</div></section>")
        html_doc = f"<!doctype html><html><head><meta charset='utf-8'></head><body>{''.join(sheets_html)}</body></html>"
        return sanitize_html(html_doc), 1, 0, warnings
    except Exception as e:
        warnings.append(f"xlsx_failed:{e}")
        return wrap_txt_to_html("XLSX parse failed"), 1, 0, warnings


def image_to_html(tmp_path: str) -> Tuple[str, int, int, list[str]]:
    if not _HAS_PIL:
        return wrap_txt_to_html("OCR not available"), 1, 0, ["ocr_tool_missing:PIL"]
    img = Image.open(tmp_path)  # type: ignore
    # light preprocessing to improve OCR
    try:
        from PIL import ImageFilter, ImageOps  # type: ignore
        img = img.convert("L")  # grayscale
        if min(img.size) < 1000:
            scale = 1000.0 / min(img.size)
            img = img.resize((int(img.width * scale), int(img.height * scale)))
        img = ImageOps.autocontrast(img)
        img = img.filter(ImageFilter.MedianFilter(size=3))
    except Exception:
        pass
    text, warnings = _ocr_image_to_text(img)
    if not text:
        warnings.append("ocr_empty")
    html_doc = wrap_txt_to_html(text or "")
    return html_doc, 1, 1 if text else 0, warnings


def html_to_html(tmp_path: str) -> Tuple[str, int, int, list[str]]:
    raw = open(tmp_path, "r", encoding="utf-8", errors="ignore").read()
    return sanitize_html(raw), 1, 0, []


def txt_to_html(tmp_path: str) -> Tuple[str, int, int, list[str]]:
    raw = open(tmp_path, "r", encoding="utf-8", errors="ignore").read()
    return wrap_txt_to_html(raw), 1, 0, []


class _BuiltInAdapter:
    """Adapter that reuses the in-house converters for generating manifests."""

    def __init__(self, converter):
        self.converter = converter

    def parse(
        self,
        file_path: str,
        mime: str,
        *,
        parse_method: str = "auto",
        prefer: Optional[str] = None,
    ):
        return self.converter(file_path, mime)


class NormalizationService:
    def __init__(self, db: DBClient, raw_store: MinioStore, canonical_store: MinioStore, *, tenant_id: str, logger, canonical_bucket: str):
        self.db = db
        self.raw_store = raw_store
        self.canonical_store = canonical_store
        self.tenant_id = tenant_id
        self.log = logger
        self.canonical_bucket = canonical_bucket
        # Prioritize AdvancedParserAdapter (MinerU/Docling) over basic BuiltIn
        adapters = [AdvancedParserAdapter(), _BuiltInAdapter(self._convert_with_builtin), SimpleFallbackAdapter()]
        self.parser_manager = DocumentParserManager(adapters=adapters)

    def _convert_with_builtin(self, file_path: str, mime: str) -> CanonicalManifest:
        m = (mime or "").lower()
        warnings: list[str] = []

        tool_name = "builtin"
        tool_version = "1.0"
        html_doc = ""
        page_count = 0
        ocr_pages = 0

        try:
            if "pdf" in m:
                tool_name, tool_version = "pymupdf", getattr(fitz, "__doc__", None) or "pymupdf" if _HAS_FITZ else "pymupdf"
                html_doc, page_count, ocr_pages, w = pdf_to_html(file_path)
                warnings.extend(w)
            elif "word" in m or m.endswith("officedocument.wordprocessingml.document") or "msword" in m or m.endswith(".docx"):
                tool_name, tool_version = "mammoth", str(getattr(mammoth, "__version__", "unknown")) if _HAS_MAMMOTH else "unavailable"
                html_doc, page_count, ocr_pages, w = docx_to_html(file_path)
                warnings.extend(w)
            elif m.endswith("officedocument.presentationml.presentation") or ":pptx" in m:
                tool_name, tool_version = "python-pptx", "py" if _HAS_PPTX else "unavailable"
                html_doc, page_count, ocr_pages, w = pptx_to_html(file_path)
                warnings.extend(w)
            elif m.endswith("officedocument.spreadsheetml.sheet") or ":xlsx" in m:
                tool_name, tool_version = "openpyxl", "py" if _HAS_OPENPYXL else "unavailable"
                html_doc, page_count, ocr_pages, w = xlsx_to_html(file_path)
                warnings.extend(w)
            elif m.startswith("image/"):
                tool_name, tool_version = ("paddleocr", "py") if _HAS_PADDLE else ("paddleocr", "unavailable")
                html_doc, page_count, ocr_pages, w = image_to_html(file_path)
                warnings.extend(w)
            elif "html" in m:
                tool_name, tool_version = "sanitize", "bs4-lxml"
                html_doc, page_count, ocr_pages, w = html_to_html(file_path)
                warnings.extend(w)
            elif m.startswith("text/") or "plain" in m:
                tool_name, tool_version = "txtwrap", "1.0"
                html_doc, page_count, ocr_pages, w = txt_to_html(file_path)
                warnings.extend(w)
            else:
                tool_name, tool_version = "fallback-txt", "1.0"
                html_doc, page_count, ocr_pages, w = txt_to_html(file_path)
                warnings.extend(["unknown_mime"] + w)
        except Exception as exc:
            warnings.append(f"builtin_parser_failed:{exc}")
            html_doc, page_count, ocr_pages, _ = txt_to_html(file_path)

        html_doc = sanitize_html(html_doc or "")
        try:
            soup2 = BeautifulSoup(html_doc, "lxml")
            changed = False
            for sec in soup2.find_all("section"):
                pg = None
                try:
                    pg = int(sec.get("data-page")) if sec.has_attr("data-page") else None
                except Exception:
                    pg = None
                if pg is not None:
                    sec["id"] = f"p-{pg}"
                    changed = True
            if changed:
                html_doc = str(soup2)
        except Exception:
            pass
        artifacts, annotated_html, stats = self._html_to_artifacts(html_doc)
        html_doc = annotated_html or html_doc
        page_count = page_count or (max((a.page_idx or 0) for a in artifacts) + 1 if artifacts else 0)

        manifest = CanonicalManifest(
            html=html_doc,
            tool_name=tool_name,
            tool_version=str(tool_version),
            page_count=page_count,
            ocr_pages=ocr_pages,
            artifacts=artifacts,
            warnings=warnings,
            stats=stats,
        )
        return manifest

    @staticmethod
    def _resolve_page_index(section: Tag, fallback: int) -> int:
        try:
            if section.has_attr("data-page"):
                return int(section.get("data-page"))
        except Exception:
            pass
        return fallback

    @staticmethod
    def _table_to_text(table: Tag) -> Tuple[str, int, int]:
        rows_text: List[str] = []
        max_cols = 0
        row_count = 0
        for tr in table.find_all("tr"):
            cells = [c.get_text(" ", strip=True) for c in tr.find_all(["th", "td"])]
            if not any(cells):
                continue
            max_cols = max(max_cols, len(cells))
            rows_text.append(" | ".join(cells))
            row_count += 1
        return "\n".join(rows_text).strip(), row_count, max_cols

    def _html_to_artifacts(self, html_doc: str) -> Tuple[List[CanonicalArtifact], str, Dict[str, Any]]:
        soup = BeautifulSoup(html_doc, "lxml")
        sections = soup.find_all("section", attrs={"data-page": True})
        if not sections:
            sections = [soup.body or soup]

        artifacts: List[CanonicalArtifact] = []
        processed_nodes: set[int] = set()
        stats: Dict[str, Any] = {
            "artifact_counts": {},
            "text_chars": 0,
            "tables": 0,
            "images": 0,
            "pages_detected": len(sections),
        }

        for idx, section in enumerate(sections):
            page_idx = self._resolve_page_index(section, idx)
            header_stack: List[str] = []

            for node in section.descendants:
                if not isinstance(node, Tag):
                    continue
                if id(node) in processed_nodes:
                    continue

                name = (node.name or "").lower()
                if name in ("script", "style", "noscript"):
                    continue

                # Skip nodes nested inside tables/lists except container tags
                if any((getattr(anc, "name", "") or "").lower() in ("table", "ul", "ol") for anc in node.parents if isinstance(anc, Tag) and anc is not section):
                    if name not in ("table", "ul", "ol"):
                        continue

                text = node.get_text(" ", strip=True)

                if name in ("h1", "h2", "h3", "h4", "h5", "h6"):
                    if not text:
                        continue
                    level = int(name[1]) if name[1:].isdigit() else 1
                    while len(header_stack) < level:
                        header_stack.append("")
                    header_stack = header_stack[: level - 1] + [text]
                    stats["artifact_counts"]["header"] = stats["artifact_counts"].get("header", 0) + 1
                    stats["text_chars"] += len(text)
                    artifact = CanonicalArtifact(
                        artifact_id=new_artifact_id("hdr"),
                        type="header",
                        text=text,
                        page_idx=page_idx,
                        headers=header_stack[:-1],
                        metadata={"level": level},
                    )
                    try:
                        node["data-artifact-id"] = artifact.artifact_id
                        node["id"] = f"a-{artifact.artifact_id}"
                    except Exception:
                        pass
                    artifacts.append(artifact)
                    processed_nodes.add(id(node))
                    continue

                if name == "p":
                    if not text:
                        continue
                    stats["artifact_counts"]["paragraph"] = stats["artifact_counts"].get("paragraph", 0) + 1
                    stats["text_chars"] += len(text)
                    artifact = CanonicalArtifact(
                        artifact_id=new_artifact_id("blk"),
                        type="paragraph",
                        text=text,
                        page_idx=page_idx,
                        headers=list(header_stack),
                    )
                    try:
                        node["data-artifact-id"] = artifact.artifact_id
                        node["id"] = f"a-{artifact.artifact_id}"
                    except Exception:
                        pass
                    artifacts.append(artifact)
                    processed_nodes.add(id(node))
                    continue

                if name in ("ul", "ol"):
                    items = [li.get_text(" ", strip=True) for li in node.find_all("li", recursive=False)]
                    items = [i for i in items if i]
                    if not items:
                        continue
                    stats["artifact_counts"]["list"] = stats["artifact_counts"].get("list", 0) + 1
                    joined = "\n".join(items)
                    stats["text_chars"] += len(joined)
                    artifact = CanonicalArtifact(
                        artifact_id=new_artifact_id("lst"),
                        type="list",
                        text=joined,
                        page_idx=page_idx,
                        headers=list(header_stack),
                        metadata={"items": len(items), "ordered": name == "ol"},
                    )
                    try:
                        node["data-artifact-id"] = artifact.artifact_id
                        node["id"] = f"a-{artifact.artifact_id}"
                    except Exception:
                        pass
                    artifacts.append(artifact)
                    processed_nodes.add(id(node))
                    continue

                if name == "table":
                    table_text, rows, cols = self._table_to_text(node)
                    if not table_text:
                        continue
                    stats["artifact_counts"]["table"] = stats["artifact_counts"].get("table", 0) + 1
                    stats["tables"] += 1
                    stats["text_chars"] += len(table_text)
                    meta = {
                        "rows": rows,
                        "cols": cols,
                        "headers": list(header_stack),
                        "html": str(node),
                    }
                    artifact = CanonicalArtifact(
                        artifact_id=new_artifact_id("tbl"),
                        type="table",
                        text=table_text,
                        page_idx=page_idx,
                        headers=list(header_stack),
                        metadata=meta,
                    )
                    try:
                        node["data-artifact-id"] = artifact.artifact_id
                        node["id"] = f"a-{artifact.artifact_id}"
                    except Exception:
                        pass
                    artifacts.append(artifact)
                    processed_nodes.add(id(node))
                    continue

                if name == "pre":
                    raw = node.get_text("\n", strip=True)
                    if not raw:
                        continue
                    stats["artifact_counts"]["code"] = stats["artifact_counts"].get("code", 0) + 1
                    stats["text_chars"] += len(raw)
                    artifact = CanonicalArtifact(
                        artifact_id=new_artifact_id("code"),
                        type="code",
                        text=raw,
                        page_idx=page_idx,
                        headers=list(header_stack),
                    )
                    try:
                        node["data-artifact-id"] = artifact.artifact_id
                        node["id"] = f"a-{artifact.artifact_id}"
                    except Exception:
                        pass
                    artifacts.append(artifact)
                    processed_nodes.add(id(node))
                    continue

                if name == "figure":
                    img = node.find("img")
                    if not img:
                        continue
                    caption_tag = node.find("figcaption")
                    caption = caption_tag.get_text(" ", strip=True) if caption_tag else None
                    alt_text = img.get("alt")
                    src = img.get("src")
                    stats["artifact_counts"]["image"] = stats["artifact_counts"].get("image", 0) + 1
                    stats["images"] += 1
                    artifact = CanonicalArtifact(
                        artifact_id=new_artifact_id("img"),
                        type="image",
                        text=(alt_text or caption or "").strip(),
                        page_idx=page_idx,
                        headers=list(header_stack),
                        caption=caption,
                        metadata={"alt": alt_text, "source": src},
                        raw_path=src,
                    )
                    try:
                        node["data-artifact-id"] = artifact.artifact_id
                        node["id"] = f"a-{artifact.artifact_id}"
                    except Exception:
                        pass
                    artifacts.append(artifact)
                    processed_nodes.add(id(node))
                    processed_nodes.add(id(img))
                    continue

                if name == "img":
                    if any(isinstance(parent, Tag) and parent.name == "figure" for parent in node.parents):
                        processed_nodes.add(id(node))
                        continue
                    alt_text = node.get("alt")
                    src = node.get("src")
                    stats["artifact_counts"]["image"] = stats["artifact_counts"].get("image", 0) + 1
                    stats["images"] += 1
                    artifact = CanonicalArtifact(
                        artifact_id=new_artifact_id("img"),
                        type="image",
                        text=(alt_text or "").strip(),
                        page_idx=page_idx,
                        headers=list(header_stack),
                        metadata={"alt": alt_text, "source": src},
                        raw_path=src,
                    )
                    try:
                        node["data-artifact-id"] = artifact.artifact_id
                        node["id"] = f"a-{artifact.artifact_id}"
                    except Exception:
                        pass
                    artifacts.append(artifact)
                    processed_nodes.add(id(node))
                    continue

        stats["artifact_total"] = len(artifacts)
        detected_pages = stats.get("pages_detected")
        if not detected_pages:
            detected_pages = max((a.page_idx or 0) for a in artifacts) + 1 if artifacts else 0
        stats["page_count_detected"] = detected_pages
        return artifacts, str(soup), stats

    @staticmethod
    def _detect_language(sample: str) -> Optional[str]:
        if not _HAS_LANG or not sample:
            return None
        try:
            lang_code = lang_detect(sample)
            return lang_code
        except Exception:
            return None

    def _blob_key_for_sha(self, sha256_hex: str) -> str:
        return self.raw_store.build_key_for_sha256(sha256_hex)

    def run_one(self, *, doc_id: str, sha256: str, mime: str) -> dict:
        t0 = time.time()
        key = self._blob_key_for_sha(sha256)

        tmp_path = self.raw_store.fget_to_tmp(key)
        try:
            parser_options = ParserOptions(
                parse_method=settings.parse_method,
                auto_ocr_fallback=settings.parse_auto_ocr_fallback,
                sparse_text_threshold=settings.parse_sparse_text_threshold,
            )
            manifest = self.parser_manager.parse(tmp_path, mime, options=parser_options)

            # Optional language detection + integrity warnings
            sample_text = " ".join(a.text for a in manifest.artifacts[:20])
            lang_code = self._detect_language(sample_text[:4000])
            if lang_code:
                manifest.warnings.append(f"lang:{lang_code}")
            if not manifest.artifacts:
                manifest.warnings.append("no_artifacts_detected")
            text_for_lang = BeautifulSoup(manifest.html or "", "lxml").get_text()
            if not (text_for_lang or "").strip():
                manifest.warnings.append("canonical_empty")

            stats = manifest.stats or {}
            if manifest.page_count:
                stats.setdefault(
                    "text_density_chars_per_page",
                    round(stats.get("text_chars", 0) / max(1, manifest.page_count), 2),
                )
            stats.setdefault("artifact_total", len(manifest.artifacts))
            stats.setdefault("ocr_pages", manifest.ocr_pages)
            if stats.get("text_chars", 0) < 200:
                manifest.warnings.append("low_text_coverage")
            if stats.get("tables", 0) and stats.get("text_chars", 0) < 100:
                manifest.warnings.append("sparse_text")
            if not manifest.stats:
                manifest.stats = stats

            # Upload canonical HTML
            canonical_key = self.canonical_store.put_canonical_html(
                bucket=self.canonical_bucket, doc_id=doc_id, html=manifest.html, version="v1"
            )

            # Persist manifest JSON (without duplicating the HTML body)
            manifest_payload = manifest.to_dict(include_html=False)
            manifest_payload.update(
                {
                    "artifact_count": len(manifest.artifacts),
                    "html_uri": canonical_key,
                }
            )
            manifest_key = self.canonical_store.put_canonical_json(
                bucket=self.canonical_bucket,
                doc_id=doc_id,
                name="manifest.json",
                payload=manifest_payload,
                version="v1",
            )

            warnings = list(dict.fromkeys(manifest.warnings))
            stats["warnings"] = warnings

            # Persist warnings back into manifest stats for downstream consumers
            manifest_payload["warnings"] = warnings
            manifest_payload["stats"] = manifest.stats

            # Persist normalization row
            self.db.insert_normalization(
                doc_id=doc_id,
                canonical_uri=canonical_key,
                tool_name=manifest.tool_name,
                tool_version=str(manifest.tool_version),
                page_count=manifest.page_count,
                ocr_pages=manifest.ocr_pages,
                warnings=warnings,
                manifest_uri=manifest_key,
            )
            self.db.update_document_state(doc_id, "NORMALIZED", ts_column="normalized_at")

            status = "OK" if not warnings else "WARN"
            self.db.insert_event(
                self.tenant_id,
                stage="NORMALIZED",
                status=status,
                details={
                    "event": "DOC_NORMALIZED_" + status,
                    "canonical_uri": canonical_key,
                    "manifest_uri": manifest_key,
                    "tool_name": manifest.tool_name,
                    "tool_version": str(manifest.tool_version),
                    "page_count": manifest.page_count,
                    "ocr_pages": manifest.ocr_pages,
                    "artifact_count": len(manifest.artifacts),
                    "warnings": warnings,
                    "stats": stats.copy(),
                },
                doc_id=doc_id,
            )

            self.log(
                "info",
                "normalized",
                stage="NORMALIZED",
                doc_id=doc_id,
                tool=manifest.tool_name,
                warnings=warnings,
                artifact_count=len(manifest.artifacts),
                stats=stats,
                latency_ms=int((time.time() - t0) * 1000),
            )

            return {
                "doc_id": doc_id,
                "status": status,
                "warnings": warnings,
                "canonical_uri": canonical_key,
                "manifest_uri": manifest_key,
                "artifact_count": len(manifest.artifacts),
            }

        except Exception as e:
            self.db.insert_event(self.tenant_id, stage="NORMALIZED", status="FAIL", details={
                "event": "DOC_NORMALIZED_FAIL", "error": str(e)
            }, doc_id=doc_id)
            self.log("error", "normalize-fail", stage="NORMALIZED", doc_id=doc_id, error=str(e))
            raise
        finally:
            try:
                os.remove(tmp_path)
            except Exception:
                pass